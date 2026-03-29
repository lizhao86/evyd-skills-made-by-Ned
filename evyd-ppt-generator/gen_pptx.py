"""
EVYD PPT Generator
==================
Renders a content.json file into a native-PPTX presentation using the EVYD
Aptos template.  All slide content is created as real shapes/textboxes —
fully editable in PowerPoint (not screenshots).

Usage:
  python gen_pptx.py content.json [--output path.pptx] [--style evyd_blue]
                                   [--template /path/to/template.pptx]

The default template path is stored in DEFAULT_TEMPLATE below.
Override per-project via content.json meta.template.
"""

import json, sys, argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from styles import STYLES

# ── Template layout indices (EVYD PPT Template Aptos.pptx) ──────────────────
L_COVER = 0   # Front Cover_Brunei
L_AGEND = 7   # Agenda_P1
L_TRANS = 29  # Transition Template_S2  (section dividers)
L_BG_B  = 31  # Background_P1  (blue content slides)
L_BG_W  = 37  # Background_White
L_END   = 39  # End_P1

DEFAULT_TEMPLATE = '/Users/Li.ZHAO/Documents/EVYD PPT Template Aptos.pptx'

I = Inches
P = Pt
NS_REL = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


# ─────────────────────────────────────────────────────────────────────────────
# Primitive helpers
# ─────────────────────────────────────────────────────────────────────────────

def _rgb(v):
    """Accept RGBColor | [r,g,b] list | '#RRGGBB' string."""
    if isinstance(v, RGBColor): return v
    if isinstance(v, (list, tuple)): return RGBColor(*v)
    if isinstance(v, str): h = v.lstrip('#'); return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
    raise ValueError(f"Cannot parse color: {v!r}")

def ph(slide, idx, text, sz=None, bold=None, color=None, align=None, italic=False, font='Aptos'):
    """Fill a placeholder by idx."""
    try: pl = slide.placeholders[idx]
    except KeyError: return
    tf = pl.text_frame; tf.clear()
    pg = tf.paragraphs[0]
    if align: pg.alignment = align
    r = pg.add_run(); r.text = text; r.font.name = font
    if sz:   r.font.size  = P(sz)
    if bold is not None: r.font.bold = bold
    if italic: r.font.italic = italic
    if color: r.font.color.rgb = _rgb(color)

def bx(slide, l, t, w, h, text='', sz=16, bold=False, color=None,
       align=None, italic=False, font='Aptos'):
    """Add a textbox. Returns the textbox shape."""
    tb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = tb.text_frame; tf.word_wrap = True
    pg = tf.paragraphs[0]
    pg.alignment = align or PP_ALIGN.LEFT
    r = pg.add_run(); r.text = text; r.font.name = font
    r.font.size = P(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = _rgb(color)
    return tb

def ap(tf, text, sz=14, bold=False, color=None, align=None, spb=0, italic=False, font='Aptos'):
    """Append a paragraph to an existing text frame."""
    pg = tf.add_paragraph()
    pg.alignment = align or PP_ALIGN.LEFT
    if spb: pg.space_before = P(spb)
    r = pg.add_run(); r.text = text; r.font.name = font
    r.font.size = P(sz); r.font.bold = bold; r.font.italic = italic
    if color: r.font.color.rgb = _rgb(color)

def rc(slide, l, t, w, h, fill=None, line=None, lw=0.75, rd=False):
    """Add a rectangle (rd=True for rounded). Returns the shape."""
    s = slide.shapes.add_shape(5 if rd else 1, I(l), I(t), I(w), I(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = _rgb(fill)
    else:    s.fill.background()
    if line: s.line.color.rgb = _rgb(line); s.line.width = P(lw)
    else:    s.line.fill.background()
    return s

def hdr(slide, section, num, title, blue=True, tsz=26, st=None):
    """Standard content-slide header: section label + slide# + title + teal rule.
    Returns content_top_y (inches) — start your content below this."""
    F = st['font'] if st else 'Aptos'
    lc = st['accent']   if st else RGBColor(0x2C,0xD5,0xC3)
    nc = st['text_num'] if st else (RGBColor(0xCC,0xE8,0xF5) if blue else RGBColor(0xAA,0xAA,0xAA))
    tc = st['white']    if st else (RGBColor(0xFF,0xFF,0xFF) if blue else RGBColor(0x17,0x2E,0x41))
    if not blue and st: tc = st['text_dark']
    if not blue and st: lc = st['accent2']
    bx(slide, 1.0, 0.30, 10,   0.28, section, sz=8,   color=lc,  font=F)
    bx(slide, 16.5, 0.30, 2.5, 0.28, num,     sz=8,   color=nc,  align=PP_ALIGN.RIGHT, font=F)
    bx(slide, 1.0, 0.70, 17,   0.75, title,   sz=tsz, bold=True, color=tc, font=F)
    rc(slide, 1.0, 1.44, 0.55, 0.045, fill=lc)
    return 1.68


# ─────────────────────────────────────────────────────────────────────────────
# Slide renderers
# ─────────────────────────────────────────────────────────────────────────────

def render_cover(slide, data, st, _num, _total):
    """
    JSON fields: title, subtitle
    Layout: L_COVER (Front Cover_Brunei)
    """
    F = st['font']
    ph(slide, 0,  data.get('title','Presentation Title'),   sz=44, bold=True, color=st['white'], font=F)
    ph(slide, 10, data.get('subtitle','Subtitle'),           sz=18, color=st['text_num'],  font=F)


def render_agenda(slide, data, st, _num, _total):
    """
    JSON fields: items: [{num, title, time}]  (max 5)
    Layout: L_AGEND (Agenda_P1)
    """
    F = st['font']
    ROWS = [(19,10),(20,11),(21,12),(22,13),(23,14)]
    for i, row in enumerate(data.get('items', [])):
        if i >= len(ROWS): break
        ci, ti = ROWS[i]
        ph(slide, ci, str(row.get('num', i+1)), sz=16, bold=True, color=st['white'], align=PP_ALIGN.CENTER, font=F)
        ph(slide, ti, f"{row['title']}   {row.get('time','')}", sz=16, color=st['white'], font=F)
    # Clear unused rows
    for ci, ti in [(25,15),(26,16),(27,17),(28,18)]:
        ph(slide, ci, ' ', font=F); ph(slide, ti, ' ', font=F)


def render_section_divider(slide, data, st, _num, _total):
    """
    JSON fields: num (e.g. "01"), title, bg_color (optional [r,g,b])
    Layout: L_TRANS (Transition_S2) — picture placeholder filled with solid rectangle
    """
    F = st['font']
    bg = _rgb(data.get('bg_color', [4, 30, 20]))
    rc(slide, 0, 0, 20, 8.48, fill=bg)
    ph(slide, 19, str(data.get('num','01')), sz=20, bold=True, color=st['white'], align=PP_ALIGN.CENTER, font=F)
    ph(slide, 20, data.get('title','Section'), sz=22, bold=True, color=st['white'], font=F)


def render_bullets_with_panel(slide, data, st, num, total):
    """
    JSON fields: section, title, background, bullets: [str],
                 side_panel: {type:"quote", text:str}
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    bullets = data.get('bullets', [])
    by = ct + 0.1
    for b in bullets:
        rc(slide, 1.0, by+0.07, 0.12, 0.12, fill=st['accent'])
        bx(slide, 1.28, by, 9.0, 0.55, b, sz=14, color=st['white'] if blue else st['text_dark'], font=F)
        by += 0.65

    # Side panel
    panel = data.get('side_panel', {})
    if panel.get('type') == 'quote':
        rc(slide, 11.0, ct+0.1, 7.9, 7.65, fill=st['card_side'])
        rc(slide, 11.0, ct+0.1, 0.12, 7.65, fill=st['accent'])
        bx(slide, 11.35, ct+0.4, 7.3, 1.1, '"', sz=80, italic=True,
           color=_rgb([0x33,0x88,0xBB]), font=F)
        bx(slide, 11.35, ct+1.5, 7.3, 4.8,
           panel.get('text',''), sz=22, italic=True, color=st['white'], font=F)

    # Optional ground rules block
    rules = data.get('ground_rules', [])
    if rules:
        rc(slide, 1.0, by+0.12, 9.5, 2.55, fill=st['card'])
        bx(slide, 1.15, by+0.20, 4, 0.28, 'GROUND RULES', sz=8, bold=True, color=st['accent'], font=F)
        ry = by + 0.58
        for rt in rules:
            bx(slide, 1.18, ry, 0.28, 0.30, '—', sz=13, bold=True, color=st['accent'], font=F)
            bx(slide, 1.50, ry, 8.8, 0.40, rt, sz=13, color=st['white'], font=F)
            ry += 0.48


def render_two_column_check(slide, data, st, num, total):
    """
    JSON fields: section, title, subtitle, background,
                 left: {title, color:[r,g,b], marker:"✓", items:[str]},
                 right: {title, color:[r,g,b], marker:"✗", items:[str]}
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct-0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_gray'], font=F)

    LIGHT_COLORS = {
        'green': [0xEB,0xF9,0xF0], 'red': [0xFF,0xEB,0xEB],
        'blue':  [0xEF,0xF7,0xFD], 'yellow':[0xFF,0xFB,0xE6],
    }
    for col_i, side in enumerate([data.get('left',{}), data.get('right',{})]):
        x = 1.0 + col_i * 9.5
        hc = _rgb(side.get('color', [0x27,0xAE,0x60] if col_i==0 else [0xE7,0x4C,0x3C]))
        lc_key = side.get('light_color','green' if col_i==0 else 'red')
        lc = _rgb(LIGHT_COLORS.get(lc_key, LIGHT_COLORS['green']))
        marker = side.get('marker', '✓' if col_i==0 else '✗')
        rc(slide, x, ct+0.3, 8.8, 7.85, fill=lc)
        rc(slide, x, ct+0.3, 8.8, 0.50, fill=hc)
        bx(slide, x+0.15, ct+0.34, 8.3, 0.38, side.get('title',''),
           sz=11, bold=True, color=st['white'], font=F)
        iy = ct + 0.98
        for item in side.get('items', []):
            bx(slide, x+0.10, iy, 0.32, 0.30, marker, sz=12, bold=True, color=hc, font=F)
            bx(slide, x+0.46, iy, 8.18, 0.45, item, sz=13.5,
               color=st['text_dark'] if not blue else st['white'], font=F)
            iy += 0.53


def render_cards_grid(slide, data, st, num, total):
    """
    JSON fields: section, title, background,
                 cards: [{num:"01", text:"..."}]  (max 6, rendered 2 per row)
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    cards = data.get('cards', [])
    CW=8.7; CH=2.55; GX=0.5; GY=0.28
    for i, card in enumerate(cards[:6]):
        col = i % 2; row = i // 2
        x = 1.0 + col*(CW+GX); y = ct + 0.1 + row*(CH+GY)
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, y, CW, CH, fill=bg)
        bx(slide, x+0.2, y+0.14, 1.5, 0.9, card.get('num',''), sz=32, bold=True,
           color=st['card_num'], font=F)
        bx(slide, x+0.2, y+1.02, CW-0.4, 1.35, card.get('text',''), sz=14,
           color=st['white'] if blue else st['text_dark'], font=F)


def render_criteria_rows(slide, data, st, num, total):
    """
    JSON fields: section, title, subtitle, background, footnote,
                 criteria: [{num:"01", label:"LABEL", text:"..."}]  (max 4)
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct-0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    rows = data.get('criteria', [])
    RH=1.85; RG=0.24
    for i, row in enumerate(rows[:4]):
        y = ct + 0.4 + i*(RH+RG)
        bg = st['card'] if blue else st['card_white']
        rc(slide, 1.0, y, 17.5, RH, fill=bg)
        bx(slide, 1.1, y+0.35, 1.3, 1.0, row.get('num',''), sz=30, bold=True,
           color=st['card_num'], font=F)
        rc(slide, 2.5, y+0.28, 0.035, RH-0.56, fill=st['line_gray'])
        bx(slide, 2.65, y+0.28, 5.0, 0.34, row.get('label',''), sz=9, bold=True,
           color=st['accent'], font=F)
        bx(slide, 2.65, y+0.70, 14.5, 0.85, row.get('text',''), sz=15,
           color=st['white'] if blue else st['text_dark'], font=F)

    if data.get('footnote'):
        fn_y = ct + 0.4 + len(rows)*(RH+RG) - 0.15
        bx(slide, 1.0, fn_y, 17.5, 0.28, data['footnote'], sz=10,
           color=st['text_dim'] if blue else st['text_gray'],
           align=PP_ALIGN.RIGHT, font=F)


def render_scope_tiers(slide, data, st, num, total):
    """
    JSON fields: section, title, subtitle, background,
                 tiers: [{color:[r,g,b], icon:"🟢", label:"...", desc:"...", example:"..."}]
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct-0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    tiers = data.get('tiers', [])
    BH=1.85; BG=0.20
    for i, tier in enumerate(tiers[:4]):
        y = ct + 0.35 + i*(BH+BG)
        tc = _rgb(tier.get('color', [0x44,0x99,0x44]))
        bg = st['card'] if blue else st['card_white']
        rc(slide, 1.0, y, 17.5, BH, fill=bg)
        rc(slide, 1.0, y, 0.14, BH, fill=tc)
        bx(slide, 1.28, y+0.15, 1.0, 0.80, tier.get('icon',''), sz=22, color=st['white'], font=F)
        bx(slide, 2.45, y+0.12, 15.3, 0.36, tier.get('label',''), sz=9, bold=True,
           color=st['accent'], font=F)
        bx(slide, 2.45, y+0.52, 15.3, 0.48, tier.get('desc',''), sz=13.5,
           color=st['white'] if blue else st['text_dark'], font=F)
        bx(slide, 2.45, y+1.08, 15.3, 0.55, tier.get('example',''), sz=11, italic=True,
           color=st['text_dim'] if blue else st['text_gray'], font=F)


def render_two_panel(slide, data, st, num, total):
    """
    JSON fields: section, title, background,
                 panels: [{color:[r,g,b], icon:"⚠️", title:"...", items:[str]}]  (2 panels)
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    for col_i, panel in enumerate(data.get('panels', [])[:2]):
        x = 1.0 + col_i * 9.5
        hc = _rgb(panel.get('color', [0x00,0x76,0xB3]))
        bg = st['card'] if blue else st['card_white']
        header_bg = st['navy']
        rc(slide, x, ct+0.25, 8.8, 7.65, fill=bg)
        rc(slide, x, ct+0.25, 8.8, 0.52, fill=header_bg)
        bx(slide, x+0.15, ct+0.31, 8.3, 0.36,
           f"{panel.get('icon','')}  {panel.get('title','')}", sz=11, bold=True,
           color=st['white'], font=F)
        iy = ct + 0.97
        for item in panel.get('items', []):
            bx(slide, x+0.10, iy, 0.32, 0.30, '▸', sz=14, bold=True, color=hc, font=F)
            bx(slide, x+0.48, iy, 8.15, 0.65, item, sz=13.5,
               color=st['white'] if blue else st['text_dark'], font=F)
            iy += 0.80


def render_two_column_steps(slide, data, st, num, total):
    """
    JSON fields: section, title, background,
                 columns: [{title:"...", steps:[{bold:"...", normal:"..."}]}]  (2 columns)
                 warning: str (optional — shown as a red bar at the bottom)
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    for col_i, col in enumerate(data.get('columns', [])[:2]):
        x = 1.0 + col_i * 9.5
        bx(slide, x, ct+0.08, 9, 0.32, col.get('title',''), sz=11, bold=True,
           color=st['accent'], font=F)
        sy = ct + 0.55
        for n, step in enumerate(col.get('steps', [])[:4], 1):
            bg = st['card'] if blue else st['card_white']
            rc(slide, x, sy, 8.8, 1.55, fill=bg)
            rc(slide, x, sy, 0.5, 1.55, fill=st['card_side'])
            bx(slide, x+0.08, sy+0.53, 0.34, 0.50, str(n), sz=14, bold=True,
               color=st['white'], align=PP_ALIGN.CENTER, font=F)
            tb = bx(slide, x+0.65, sy+0.22, 8.0, 0.52,
                    step.get('bold',''), sz=14, bold=True, color=st['white'], font=F)
            ap(tb.text_frame, step.get('normal',''), sz=12.5,
               color=st['text_dim'] if blue else st['text_gray'], font=F)
            sy += 1.72

    if data.get('warning'):
        rc(slide, 1.0, 8.60, 17.5, 0.72, fill=_rgb([0x3A,0x08,0x08]))
        rc(slide, 1.0, 8.60, 0.14, 0.72, fill=_rgb([0xE7,0x4C,0x3C]))
        bx(slide, 1.25, 8.66, 17.0, 0.55, data['warning'], sz=12,
           color=_rgb([0xFF,0xAA,0xAA]), font=F)


def render_scenario_cards(slide, data, st, num, total):
    """
    JSON fields: section, title, subtitle, background,
                 scenarios: [{num:"01", icon:"💊", title:"...", desc:"...", tag:"..."}]
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'blue') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct-0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_dim'] if blue else st['text_gray'], font=F)

    SW=8.7; SH=3.48
    for i, scen in enumerate(data.get('scenarios', [])[:4]):
        col = i % 2; row = i // 2
        x = 1.0 + col*(SW+0.5); y = ct + 0.35 + row*(SH+0.28)
        bg = st['card'] if blue else st['card_white']
        rc(slide, x, y, SW, SH, fill=bg)
        rc(slide, x, y, SW, 0.08, fill=st['accent'])
        bx(slide, x+0.2, y+0.18, 2.5, 0.28,
           f"Scenario {scen.get('num','')}", sz=9, bold=True, color=st['accent'], font=F)
        bx(slide, x+0.2, y+0.52, SW-0.4, 0.58,
           f"{scen.get('icon','')}  {scen.get('title','')}", sz=15, bold=True,
           color=st['white'] if blue else st['text_dark'], font=F)
        bx(slide, x+0.2, y+1.18, SW-0.4, 1.88, scen.get('desc',''), sz=13,
           color=st['text_dim'] if blue else st['text_gray'], font=F)
        rc(slide, x+0.2, y+SH-0.55, 2.5, 0.38, fill=st['card_side'])
        bx(slide, x+0.26, y+SH-0.51, 2.35, 0.28, scen.get('tag',''), sz=9,
           color=st['accent'], font=F)


def render_survey(slide, data, st, num, total):
    """
    JSON fields: section, title, subtitle, background,
                 steps: [{title:"...", desc:"..."}]
                 qr_label, qr_note, qr_hint
    Layout: L_BG_B or L_BG_W
    """
    blue = data.get('background', 'white') == 'blue'
    F = st['font']
    ct = hdr(slide, data.get('section',''), _slide_num(data, num, total),
             data.get('title',''), blue=blue, st=st)

    if data.get('subtitle'):
        bx(slide, 1.0, ct-0.05, 17, 0.28, data['subtitle'], sz=12,
           color=st['text_gray'], font=F)

    STH=1.78
    for i, step in enumerate(data.get('steps', [])[:4]):
        y = ct + 0.35 + i*(STH+0.22)
        bg = st['card_white'] if not blue else st['card']
        rc(slide, 1.0, y, 10.5, STH, fill=bg)
        rc(slide, 1.0, y, 0.65, STH, fill=st['accent2'])
        bx(slide, 1.06, y+0.56, 0.53, 0.65, str(i+1), sz=22, bold=True,
           color=st['white'], align=PP_ALIGN.CENTER, font=F)
        bx(slide, 1.82, y+0.26, 9.5, 0.48, step.get('title',''), sz=15, bold=True,
           color=st['text_dark'] if not blue else st['white'], font=F)
        bx(slide, 1.82, y+0.84, 9.5, 0.72, step.get('desc',''), sz=13,
           color=st['text_gray'], font=F)

    # QR panel
    qr_bg = st['card_white'] if not blue else st['card']
    rc(slide, 12.3, ct+0.3, 6.6, 9.0, fill=qr_bg, line=st['accent2'], lw=1.0)
    bx(slide, 12.3, ct+0.45, 6.6, 0.34,
       data.get('qr_label','SURVEY QR CODE'), sz=9, bold=True,
       color=st['accent2'], align=PP_ALIGN.CENTER, font=F)
    rc(slide, 13.25, ct+1.05, 4.6, 4.8, fill=_rgb([0xD8,0xEA,0xF5]),
       line=st['accent2'], lw=0.75)
    bx(slide, 13.25, ct+2.95, 4.6, 1.0, '[ QR Code ]', sz=14, italic=True,
       color=_rgb([0x66,0x99,0xBB]), align=PP_ALIGN.CENTER, font=F)
    bx(slide, 12.3, ct+6.1, 6.6, 0.50,
       data.get('qr_note','(Link provided on the day)'), sz=11, italic=True,
       color=_rgb([0x77,0x99,0xBB]), align=PP_ALIGN.CENTER, font=F)
    bx(slide, 12.3, ct+6.8, 6.6, 0.80,
       data.get('qr_hint','Moderator will circulate\nto assist with submission'),
       sz=12, color=st['text_gray'], align=PP_ALIGN.CENTER, font=F)


def render_ending(slide, data, st, _num, _total):
    """
    JSON fields: title, subtitle, actions: [{icon, title, desc}]  (2 actions)
    Layout: L_END (End_P1) — EVYD logo provided by layout
    """
    F = st['font']
    bx(slide, 4.0, 1.3, 12.0, 0.75,
       data.get('title','Thank You'), sz=32, bold=True,
       color=st['white'], align=PP_ALIGN.CENTER, font=F)
    bx(slide, 4.0, 2.15, 12.0, 0.48,
       data.get('subtitle',''), sz=15,
       color=st['text_num'], align=PP_ALIGN.CENTER, font=F)

    for col_i, action in enumerate(data.get('actions',[])[:2]):
        x = 4.5 + col_i * 5.8
        rc(slide, x, 7.55, 5.1, 1.70, fill=_rgb([0x00,0x4A,0x85]))
        bx(slide, x+0.2, 7.65, 4.5, 0.45,
           f"{action.get('icon','')}  {action.get('title','')}", sz=13, bold=True,
           color=st['accent'], font=F)
        bx(slide, x+0.2, 8.18, 4.5, 0.90,
           action.get('desc',''), sz=12, color=st['white'], font=F)


# ─────────────────────────────────────────────────────────────────────────────
# Routing
# ─────────────────────────────────────────────────────────────────────────────

RENDERERS = {
    'cover':               (L_COVER, render_cover),
    'agenda':              (L_AGEND, render_agenda),
    'section_divider':     (L_TRANS, render_section_divider),
    'bullets_with_panel':  (None,    render_bullets_with_panel),   # layout by background
    'two_column_check':    (None,    render_two_column_check),
    'cards_grid':          (None,    render_cards_grid),
    'criteria_rows':       (None,    render_criteria_rows),
    'scope_tiers':         (None,    render_scope_tiers),
    'two_panel':           (None,    render_two_panel),
    'two_column_steps':    (None,    render_two_column_steps),
    'scenario_cards':      (None,    render_scenario_cards),
    'survey':              (None,    render_survey),
    'ending':              (L_END,   render_ending),
}

FIXED_LAYOUTS = {'cover', 'agenda', 'section_divider', 'ending'}

def _bg_layout(slide_data):
    return L_BG_W if slide_data.get('background','blue') == 'white' else L_BG_B

def _slide_num(slide_data, auto_num, total):
    return slide_data.get('num', f'{auto_num:02d} / {total:02d}')


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description='EVYD PPT Generator')
    parser.add_argument('content', help='Path to content JSON file')
    parser.add_argument('--template', default=None,
                        help=f'PPTX template path (default: {DEFAULT_TEMPLATE})')
    parser.add_argument('--output', '-o', default=None, help='Output .pptx path')
    parser.add_argument('--style', default=None,
                        help='Style preset name (default: from JSON meta.style)')
    args = parser.parse_args()

    with open(args.content, encoding='utf-8') as f:
        content = json.load(f)

    meta = content.get('meta', {})
    style_name = args.style or meta.get('style', 'evyd_blue')
    if style_name not in STYLES:
        print(f'Unknown style "{style_name}". Available: {", ".join(STYLES)}')
        sys.exit(1)
    st = STYLES[style_name]

    template_path = args.template or meta.get('template', DEFAULT_TEMPLATE)
    output_path   = args.output   or meta.get('output',   'output.pptx')

    prs = Presentation(template_path)
    n_template = len(prs.slides)   # count before we add ours

    slides = content['slides']
    total  = len(slides)
    counter = 0  # auto slide-number counter (excludes chrome slides)

    for slide_data in slides:
        stype = slide_data.get('type', '')
        if stype not in RENDERERS:
            print(f'  ⚠  Unknown slide type "{stype}" — skipped')
            continue

        fixed_layout, renderer = RENDERERS[stype]
        layout_idx = fixed_layout if stype in FIXED_LAYOUTS else _bg_layout(slide_data)
        s = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if stype not in FIXED_LAYOUTS:
            counter += 1

        renderer(s, slide_data, st, counter, total)

    # Remove the original template slides (they were loaded first)
    sldIdLst = prs.slides._sldIdLst
    removed = 0
    for _ in range(n_template):
        if len(sldIdLst) == 0: break
        first  = sldIdLst[0]
        rid    = first.get(f'{{{NS_REL}}}id')
        sldIdLst.remove(first)
        try: prs.part.rels.pop(rid)
        except: pass
        removed += 1

    prs.save(output_path)
    print(f'Removed {removed} template slides')
    print(f'Final slide count: {len(prs.slides)}')
    print(f'✓ Saved → {output_path}')


if __name__ == '__main__':
    import argparse
    main()
